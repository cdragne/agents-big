"""
Chatbot Agents BIG — application Streamlit
Dashboard avec 7 agents specialises (persona + connaissances metier),
chacun discutable via l'API Anthropic (Claude).

Fonctionnalites documents :
- Documents communs : tout fichier depose dans le dossier "documents/" du repo
  est lu automatiquement et connu par tous les agents en permanence.
- Documents ponctuels : chaque agent permet d'uploader un fichier (PDF, Word,
  Excel, texte) valable pour la conversation en cours uniquement.

Lancement local : streamlit run streamlit_app.py
Deploiement gratuit : voir README.md
"""

import os
import io
import streamlit as st
from anthropic import Anthropic
from pypdf import PdfReader
import docx
import openpyxl
import xlrd
from pptx import Presentation
from fpdf import FPDF

# ---------------------------------------------------------------------------
# Configuration generale
# ---------------------------------------------------------------------------

st.set_page_config(
    page_title="Agents BIG",
    page_icon="🏭",
    layout="wide",
)

MODEL = "claude-sonnet-4-5"

DOSSIER_DOCUMENTS_COMMUNS = "documents"
TAILLE_MAX_PAR_DOCUMENT = 12000       # caracteres max extraits par document
TAILLE_MAX_TOTALE_CONTEXTE = 40000    # caracteres max, tous documents confondus

CONTEXTE_ENTREPRISE = """
Contexte entreprise (a toujours garder en tete) :
- Bordas Industrial Group (BIG), Ternay (69360), France.
- Christophe Dragne Durret est Directeur du developpement et Directeur national maintenance.
- Activite : prestations de maintenance de dallage des sols industriels et amelioration
  des sols (et developpement de la resine de sol).
- Clients : decideurs, principalement dans la distribution et la logistique.
- Prospection : telephonique, LinkedIn, Google.
- Equipe commerciale : Marie Guerreiro, Matthis Pineau, Emeric Nieps (charges d'affaires),
  Warda (assistante), + equipes operationnelles terrain.
- Outils internes : Notion, Office 365.
- Priorite actuelle : developpement national et proche Europe, developpement de la resine
  de sol, rentabilite des affaires, automatisation des process. Objectif a 5 mois : rachat
  de l'entreprise.
- Christophe est non-technique : reponses claires, directes, sans jargon, format concis
  (texte redige ou listes courtes selon le besoin), toujours en francais.

Important : tu n'as pas d'acces direct a Notion, Outlook ou aux autres outils internes de
BIG dans cette application. Si une question necessite des donnees precises que tu n'as pas
(chiffres exacts, liste d'affaires en cours...), dis-le clairement et demande a Christophe
de te donner les chiffres ou de coller les donnees, plutot que d'inventer des valeurs.

Des documents peuvent t'etre fournis (documents communs de l'entreprise et/ou documents
ajoutes par Christophe pour cette conversation). Ils apparaissent ci-dessous entre les
balises <documents_communs> et <documents_conversation>. Appuie-toi dessus en priorite
quand ils sont pertinents pour la question posee, et cite le nom du document si tu t'en
sers. Si les documents ne contiennent pas l'information demandee, dis-le clairement.
"""

AGENTS = {
    "Commercial": {
        "icone": "📞",
        "resume": "Prospection, argumentaires, emails, suivi du pipeline.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent COMMERCIAL de BIG. Tu aides Christophe et son equipe (Marie, Matthis, Emeric)
sur :
- la preparation d'appels de prospection (fiche de preparation : contexte prospect,
  angle d'approche, objections probables, questions a poser),
- les argumentaires de vente pour le dallage et la resine de sol, adaptes au secteur et
  a l'interlocuteur,
- la redaction d'emails commerciaux (prospection, relance, suivi de devis, apres RDV),
- l'analyse et la priorisation du pipeline commercial (si Christophe te donne les donnees
  des affaires en cours, aide-le a les prioriser et a identifier les relances a faire).
Sois concret, oriente action, et pense toujours du point de vue d'un decideur logistique
ou distribution (temps d'arret, securite, continuite d'exploitation, cout total).
""",
    },
    "Marketing": {
        "icone": "📣",
        "resume": "Contenus, supports commerciaux, veille sectorielle, performance.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent MARKETING de BIG. Tu aides Christophe sur :
- la creation de contenu (posts LinkedIn, newsletters, articles) dans la voix de Christophe,
  sans ton commercial trop appuye,
- les supports commerciaux ecrits (plaquettes, one-pagers, presentations client) adaptes
  au dallage ou a la resine et au contexte du prospect,
- la veille sectorielle (nouveaux entrepots, tendances sols industriels, actualites
  concurrents, mouvements dans la distribution/logistique) — precise que tu t'appuies sur
  tes connaissances generales et invite Christophe a verifier les infos recentes,
- l'analyse de performance marketing si des donnees (LinkedIn, resultats de posts) sont
  fournies par Christophe.
""",
    },
    "Manager": {
        "icone": "🧭",
        "resume": "Pilotage d'equipe, reunions, alertes, reporting KPI.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent MANAGER de BIG. Tu aides Christophe a piloter son equipe commerciale
(Marie Guerreiro, Matthis Pineau, Emeric Nieps) et son activite :
- preparer des points d'equipe et reunions (ordre du jour structure),
- detecter les signaux de derive ou de risque si Christophe partage des donnees
  d'activite ou de pipeline,
- structurer un reporting KPI (commercial, hebdomadaire ou mensuel) a partir des chiffres
  que Christophe te donne,
- proposer des actions manageriales concretes (relance, coaching, arbitrage de priorites).
Ne fabrique jamais de chiffres : si tu n'as pas les donnees, demande-les.
""",
    },
    "Comptable / Financier": {
        "icone": "📊",
        "resume": "Analyse financiere, budget, rentabilite, tresorerie.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent COMPTABLE / FINANCIER de BIG. Tu aides Christophe sur :
- l'analyse financiere globale (compte de resultat, EBE, ratios, pistes de valorisation
  en vue d'une cession/rachat — objectif de Christophe a 5 mois),
- le suivi budgetaire (budget vs realise, ecarts, projection d'atterrissage),
- le controle de rentabilite d'une affaire (marge, devis vs realise, cout de revient),
- le suivi de tresorerie (flux, delais de paiement, DSO, retards clients).
Base-toi uniquement sur les chiffres fournis par Christophe. Presente toujours les
hypotheses et les limites de ton calcul. Rappelle, si pertinent, que tu ne remplaces pas
un expert-comptable pour les decisions officielles.
""",
    },
    "Juridique": {
        "icone": "⚖️",
        "resume": "Contrats, CGV, litiges, cession, RGPD.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent JURIDIQUE de BIG. Tu aides Christophe a comprendre les enjeux juridiques
lies a son activite : contrats, CGV, sous-traitance, impayes/mise en demeure, droit du
travail, appels d'offres, RGPD, et les sujets lies au projet de cession/rachat de
l'entreprise (due diligence, pacte d'associes).
Donne une explication claire, en francais simple, des risques et options — et rappelle
systematiquement que tu ne remplaces pas un avocat : pour toute decision a enjeu, une
verification par un professionnel du droit est necessaire.
""",
    },
    "Technique sols industriels": {
        "icone": "🏗️",
        "resume": "Dallage, resine, pathologies, produits, formations, memoires techniques.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent TECHNIQUE sols industriels de BIG. Tu aides Christophe et ses equipes sur :
- les questions techniques sur le dallage beton et les revetements resine (pathologies,
  fissures, planeite, normes, diagnostic),
- les recommandations de produits (resine, ragreage, primaire d'accrochage, produits ESD
  ou alimentaires) — en indiquant qu'il faut confirmer avec les fiches techniques
  fabricant reelles avant mise en oeuvre,
- les formations/habilitations necessaires pour les interventions (CACES, certifications
  applicateur, securite chantier),
- la redaction de memoires techniques et reponses a appels d'offres.
Reste factuel et precis ; signale quand une verification terrain ou un avis d'expert
est indispensable avant toute preconisation engageante.
""",
    },
    "Administratif": {
        "icone": "🗂️",
        "resume": "Emails, taches, priorites du jour.",
        "system": CONTEXTE_ENTREPRISE + """
Tu es l'agent ADMINISTRATIF de BIG. Tu aides Christophe a :
- structurer et prioriser ses emails (si Christophe copie-colle des messages ou des
  resumes de sa boite mail),
- creer et organiser des listes de taches,
- faire un point rapide de priorites de journee a partir de ce que Christophe te
  partage (RDV, echeances, taches en cours).
Cette application n'a pas d'acces direct a Outlook, Notion ou Microsoft To Do : demande
a Christophe de coller le contenu pertinent si besoin.
""",
    },
}

# ---------------------------------------------------------------------------
# Extraction de texte depuis les documents
# ---------------------------------------------------------------------------


def _tronquer(texte: str, limite: int = TAILLE_MAX_PAR_DOCUMENT) -> str:
    if len(texte) > limite:
        return texte[:limite] + "\n[...document tronque, trop volumineux...]"
    return texte


def lire_pdf(donnees: bytes) -> str:
    lecteur = PdfReader(io.BytesIO(donnees))
    pages = [page.extract_text() or "" for page in lecteur.pages]
    return "\n".join(pages)


def lire_docx(donnees: bytes) -> str:
    document = docx.Document(io.BytesIO(donnees))
    paragraphes = [p.text for p in document.paragraphs]
    return "\n".join(paragraphes)


def lire_xlsx(donnees: bytes) -> str:
    classeur = openpyxl.load_workbook(io.BytesIO(donnees), data_only=True)
    morceaux = []
    for nom_feuille in classeur.sheetnames:
        feuille = classeur[nom_feuille]
        morceaux.append(f"--- Feuille : {nom_feuille} ---")
        for ligne in feuille.iter_rows(values_only=True):
            valeurs = [str(v) for v in ligne if v is not None]
            if valeurs:
                morceaux.append(" | ".join(valeurs))
    return "\n".join(morceaux)


def lire_xls(donnees: bytes) -> str:
    classeur = xlrd.open_workbook(file_contents=donnees)
    morceaux = []
    for feuille in classeur.sheets():
        morceaux.append(f"--- Feuille : {feuille.name} ---")
        for num_ligne in range(feuille.nrows):
            valeurs = [str(v) for v in feuille.row_values(num_ligne) if v not in ("", None)]
            if valeurs:
                morceaux.append(" | ".join(valeurs))
    return "\n".join(morceaux)


def lire_txt(donnees: bytes) -> str:
    return donnees.decode("utf-8", errors="ignore")


def extraire_texte(nom_fichier: str, donnees: bytes) -> str:
    extension = nom_fichier.lower().rsplit(".", 1)[-1] if "." in nom_fichier else ""
    try:
        if extension == "pdf":
            texte = lire_pdf(donnees)
        elif extension == "docx":
            texte = lire_docx(donnees)
        elif extension == "xlsx":
            texte = lire_xlsx(donnees)
        elif extension == "xls":
            texte = lire_xls(donnees)
        elif extension in ("txt", "md", "csv"):
            texte = lire_txt(donnees)
        else:
            return f"[Format .{extension} non pris en charge pour {nom_fichier}]"
    except Exception as erreur:
        return f"[Erreur de lecture de {nom_fichier} : {erreur}]"
    return _tronquer(texte)


@st.cache_data(show_spinner=False)
def charger_documents_communs():
    """Lit tous les fichiers du dossier documents/ (une seule fois, mis en cache)."""
    resultats = []
    if os.path.isdir(DOSSIER_DOCUMENTS_COMMUNS):
        for nom_fichier in sorted(os.listdir(DOSSIER_DOCUMENTS_COMMUNS)):
            chemin = os.path.join(DOSSIER_DOCUMENTS_COMMUNS, nom_fichier)
            if os.path.isfile(chemin):
                with open(chemin, "rb") as f:
                    donnees = f.read()
                texte = extraire_texte(nom_fichier, donnees)
                resultats.append({"nom": nom_fichier, "texte": texte})
    return resultats


def construire_bloc_documents(documents: list, balise: str) -> str:
    if not documents:
        return f"<{balise}>\n(aucun document)\n</{balise}>"
    morceaux = [f"<{balise}>"]
    total = 0
    for doc in documents:
        if total >= TAILLE_MAX_TOTALE_CONTEXTE:
            morceaux.append("[...limite globale de documents atteinte, documents suivants ignores...]")
            break
        morceaux.append(f"--- Document : {doc['nom']} ---\n{doc['texte']}")
        total += len(doc["texte"])
    morceaux.append(f"</{balise}>")
    return "\n".join(morceaux)


# ---------------------------------------------------------------------------
# Generation de documents (Word, Excel, PDF, PowerPoint) a partir d'une reponse
# ---------------------------------------------------------------------------

MOTS_CLES_EXPORT = [
    "word", ".docx", "excel", ".xlsx", "pdf", "powerpoint", "pptx",
    "diaporama", "presentation", "présentation", "genere un document",
    "génère un document", "genere un fichier", "génère un fichier",
    "exporte", "télécharge", "telecharge",
]


def question_demande_export(question: str) -> bool:
    q = question.lower()
    return any(mot in q for mot in MOTS_CLES_EXPORT)


def analyser_markdown_simple(texte: str):
    """Decoupe une reponse texte/markdown en elements structurés
    (titres, puces, paragraphes) reutilisables pour chaque format de sortie."""
    elements = []
    for ligne in texte.split("\n"):
        ligne_strip = ligne.strip()
        if not ligne_strip:
            continue
        if ligne_strip.startswith("### "):
            elements.append(("titre3", ligne_strip[4:].strip()))
        elif ligne_strip.startswith("## "):
            elements.append(("titre2", ligne_strip[3:].strip()))
        elif ligne_strip.startswith("# "):
            elements.append(("titre1", ligne_strip[2:].strip()))
        elif ligne_strip.startswith(("- ", "* ")):
            elements.append(("puce", ligne_strip[2:].strip()))
        else:
            elements.append(("paragraphe", ligne_strip))
    return elements


def generer_word(texte: str, titre_document: str) -> bytes:
    document = docx.Document()
    document.add_heading(titre_document, level=0)
    for type_element, contenu in analyser_markdown_simple(texte):
        if type_element == "titre1":
            document.add_heading(contenu, level=1)
        elif type_element == "titre2":
            document.add_heading(contenu, level=2)
        elif type_element == "titre3":
            document.add_heading(contenu, level=3)
        elif type_element == "puce":
            document.add_paragraph(contenu, style="List Bullet")
        else:
            document.add_paragraph(contenu)
    tampon = io.BytesIO()
    document.save(tampon)
    return tampon.getvalue()


def generer_excel(texte: str, titre_document: str) -> bytes:
    classeur = openpyxl.Workbook()
    feuille = classeur.active
    feuille.title = "Document"
    feuille["A1"] = titre_document
    feuille["A1"].font = openpyxl.styles.Font(bold=True, size=14)
    ligne_actuelle = 3
    for type_element, contenu in analyser_markdown_simple(texte):
        if type_element in ("titre1", "titre2", "titre3"):
            cellule = feuille.cell(row=ligne_actuelle, column=1, value=contenu)
            cellule.font = openpyxl.styles.Font(bold=True)
        elif " | " in contenu:
            for i, valeur in enumerate(contenu.split(" | ")):
                feuille.cell(row=ligne_actuelle, column=1 + i, value=valeur.strip())
        elif type_element == "puce":
            feuille.cell(row=ligne_actuelle, column=1, value=f"• {contenu}")
        else:
            feuille.cell(row=ligne_actuelle, column=1, value=contenu)
        ligne_actuelle += 1
    feuille.column_dimensions["A"].width = 80
    tampon = io.BytesIO()
    classeur.save(tampon)
    return tampon.getvalue()


def generer_pdf(texte: str, titre_document: str) -> bytes:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, titre_document.encode("latin-1", "replace").decode("latin-1"))
    pdf.ln(4)
    tailles = {"titre1": 14, "titre2": 12, "titre3": 11}
    for type_element, contenu in analyser_markdown_simple(texte):
        contenu_propre = contenu.encode("latin-1", "replace").decode("latin-1")
        if type_element in tailles:
            pdf.set_font("Helvetica", "B", tailles[type_element])
            pdf.multi_cell(0, 7, contenu_propre)
        elif type_element == "puce":
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(0, 6, f"-  {contenu_propre}")
        else:
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(0, 6, contenu_propre)
        pdf.ln(1)
    return bytes(pdf.output())


def generer_pptx(texte: str, titre_document: str) -> bytes:
    presentation = Presentation()

    diapo_titre = presentation.slides.add_slide(presentation.slide_layouts[0])
    diapo_titre.shapes.title.text = titre_document
    if len(diapo_titre.placeholders) > 1:
        diapo_titre.placeholders[1].text = "Genere par l'agent BIG"

    def nouvelle_diapo(titre_diapo: str):
        diapo = presentation.slides.add_slide(presentation.slide_layouts[1])
        diapo.shapes.title.text = titre_diapo
        return diapo.placeholders[1].text_frame

    zone_texte = nouvelle_diapo("Contenu")
    premier_ajout = True

    for type_element, contenu in analyser_markdown_simple(texte):
        if type_element in ("titre1", "titre2"):
            zone_texte = nouvelle_diapo(contenu)
            premier_ajout = True
        else:
            if premier_ajout:
                paragraphe = zone_texte.paragraphs[0]
                premier_ajout = False
            else:
                paragraphe = zone_texte.add_paragraph()
            paragraphe.text = contenu

    tampon = io.BytesIO()
    presentation.save(tampon)
    return tampon.getvalue()


def nom_fichier_propre(chaine: str) -> str:
    autorises = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789_-"
    nettoye = "".join(c if c in autorises else "_" for c in chaine)
    return nettoye[:60]


def afficher_boutons_export(contenu: str, agent: str, cle_unique: str):
    titre_document = f"Reponse - Agent {agent}"
    nom_base = nom_fichier_propre(f"{agent}_{cle_unique}")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.download_button(
            "📄 Word", data=generer_word(contenu, titre_document),
            file_name=f"{nom_base}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            key=f"word_{nom_base}", use_container_width=True,
        )
    with col2:
        st.download_button(
            "📊 Excel", data=generer_excel(contenu, titre_document),
            file_name=f"{nom_base}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f"excel_{nom_base}", use_container_width=True,
        )
    with col3:
        st.download_button(
            "📕 PDF", data=generer_pdf(contenu, titre_document),
            file_name=f"{nom_base}.pdf", mime="application/pdf",
            key=f"pdf_{nom_base}", use_container_width=True,
        )
    with col4:
        st.download_button(
            "📽️ PowerPoint", data=generer_pptx(contenu, titre_document),
            file_name=f"{nom_base}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"pptx_{nom_base}", use_container_width=True,
        )


# ---------------------------------------------------------------------------
# Etat de session
# ---------------------------------------------------------------------------

if "agent_actif" not in st.session_state:
    st.session_state.agent_actif = None

if "historiques" not in st.session_state:
    st.session_state.historiques = {nom: [] for nom in AGENTS}

if "documents_conversation" not in st.session_state:
    st.session_state.documents_conversation = {nom: [] for nom in AGENTS}


def get_client():
    api_key = os.environ.get("ANTHROPIC_API_KEY") or st.secrets.get("ANTHROPIC_API_KEY", None)
    if not api_key:
        st.error(
            "Aucune cle API trouvee. Ajoutez ANTHROPIC_API_KEY dans les secrets de "
            "l'application (voir README.md)."
        )
        st.stop()
    return Anthropic(api_key=api_key)


# ---------------------------------------------------------------------------
# Barre laterale : dashboard des agents
# ---------------------------------------------------------------------------

with st.sidebar:
    st.title("🏭 Agents BIG")
    st.caption("Choisissez un agent pour demarrer une conversation.")
    for nom, infos in AGENTS.items():
        bouton_label = f"{infos['icone']} {nom}"
        if st.button(bouton_label, use_container_width=True, key=f"btn_{nom}"):
            st.session_state.agent_actif = nom
    st.divider()
    if st.session_state.agent_actif:
        if st.button("🗑️ Effacer cette conversation", use_container_width=True):
            st.session_state.historiques[st.session_state.agent_actif] = []
            st.session_state.documents_conversation[st.session_state.agent_actif] = []
            st.rerun()

    st.divider()
    documents_communs = charger_documents_communs()
    with st.expander(f"📚 Documents communs ({len(documents_communs)})"):
        if documents_communs:
            for doc in documents_communs:
                st.caption(f"• {doc['nom']}")
        else:
            st.caption(
                "Aucun document commun. Deposez des fichiers dans le dossier "
                "'documents/' du repository GitHub pour les rendre disponibles "
                "a tous les agents."
            )

# ---------------------------------------------------------------------------
# Zone principale
# ---------------------------------------------------------------------------

if not st.session_state.agent_actif:
    st.title("Bienvenue sur votre assistant BIG")
    st.write(
        "Selectionnez un agent dans le menu a gauche pour commencer. "
        "Chaque agent est specialise sur un domaine de votre activite."
    )
    cols = st.columns(3)
    for i, (nom, infos) in enumerate(AGENTS.items()):
        with cols[i % 3]:
            st.subheader(f"{infos['icone']} {nom}")
            st.write(infos["resume"])
else:
    agent = st.session_state.agent_actif
    infos = AGENTS[agent]
    st.title(f"{infos['icone']} Agent {agent}")
    st.caption(infos["resume"])

    # --- Documents pour cette conversation ---
    with st.expander("📎 Documents pour cette conversation"):
        fichiers_uploades = st.file_uploader(
            "Ajouter un ou plusieurs documents (PDF, Word, Excel, texte)",
            type=["pdf", "docx", "xlsx", "xls", "txt", "md", "csv"],
            accept_multiple_files=True,
            key=f"upload_{agent}",
        )
        if fichiers_uploades:
            noms_deja_charges = {d["nom"] for d in st.session_state.documents_conversation[agent]}
            for fichier in fichiers_uploades:
                if fichier.name not in noms_deja_charges:
                    texte = extraire_texte(fichier.name, fichier.read())
                    st.session_state.documents_conversation[agent].append(
                        {"nom": fichier.name, "texte": texte}
                    )

        docs_agent = st.session_state.documents_conversation[agent]
        if docs_agent:
            st.caption("Documents actifs dans cette conversation :")
            for i, doc in enumerate(docs_agent):
                col1, col2 = st.columns([5, 1])
                col1.write(f"📄 {doc['nom']}")
                if col2.button("✕", key=f"suppr_{agent}_{i}"):
                    st.session_state.documents_conversation[agent].pop(i)
                    st.rerun()
        else:
            st.caption("Aucun document ajoute pour cette conversation.")

    for idx, message in enumerate(st.session_state.historiques[agent]):
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            if message["role"] == "assistant":
                afficher_boutons_export(message["content"], agent, f"hist{idx}")

    question = st.chat_input(f"Ecrivez a l'agent {agent}...")
    if question:
        st.session_state.historiques[agent].append({"role": "user", "content": question})
        with st.chat_message("user"):
            st.markdown(question)

        # Construction du system prompt avec documents (communs + conversation)
        bloc_communs = construire_bloc_documents(documents_communs, "documents_communs")
        bloc_conversation = construire_bloc_documents(
            st.session_state.documents_conversation[agent], "documents_conversation"
        )
        system_avec_documents = infos["system"] + "\n\n" + bloc_communs + "\n\n" + bloc_conversation

        client = get_client()
        with st.chat_message("assistant"):
            placeholder = st.empty()
            reponse_complete = ""
            with client.messages.stream(
                model=MODEL,
                max_tokens=2000,
                system=system_avec_documents,
                messages=st.session_state.historiques[agent],
            ) as stream:
                for texte in stream.text_stream:
                    reponse_complete += texte
                    placeholder.markdown(reponse_complete)

        st.session_state.historiques[agent].append(
            {"role": "assistant", "content": reponse_complete}
        )

        if question_demande_export(question):
            st.info("📎 Document pret : choisissez le format ci-dessous.")
        afficher_boutons_export(
            reponse_complete, agent, f"nouv{len(st.session_state.historiques[agent])}"
        )
